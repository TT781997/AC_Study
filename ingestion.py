"""
ingestion.py — Módulo de Ingestão Universal (Universal ScholarGPT v5.2)
========================================================================

Parser multi-formato com fallback de OCR para PDFs digitalizados.

Formatos suportados:
  • PDF                   → PyMuPDF (texto digital) + fallback OCR se < 100 chars
  • DOCX                  → python-docx
  • PPTX                  → python-pptx
  • XLSX                  → openpyxl (via pandas se disponível, senão direto)
  • TXT / MD / HTML       → leitura directa (HTML com strip de tags básico)
  • PNG / JPG / JPEG      → pipeline Vision (via callback injectada)

API pública:
    ingest_file(file, *, vision_callback=None) -> dict
        Retorna {'text': str, 'metadata': {...}, 'sources': {...}}

Design:
  • Todos os parsers em try/except → erros tornam-se warnings amigáveis
    no dict de retorno, NUNCA tracebacks crus.
  • `vision_callback` é injectada (não importa de llm_client) → testável.
    Assinatura esperada:  callable(image_bytes: bytes) -> str
  • Graceful degradation: se python-docx/pptx/openpyxl não estão instalados,
    o ficheiro é rejeitado com mensagem clara em vez de crash.
  • Fallback OCR em camadas:
      1. vision_callback (se passada)            → melhor qualidade
      2. pytesseract + pdf2image (se instalados) → gratuito, local
      3. Skip + aviso                            → último recurso
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from typing import Any, Callable, Dict, List, Optional, Tuple

from db import compute_file_hash

# ─────────────────────────────────────────────────────────────────────────────
# CONSTANTES
# ─────────────────────────────────────────────────────────────────────────────

SUPPORTED_FORMATS: List[str] = [
    "pdf", "docx", "pptx", "xlsx",
    "txt", "html", "md",
    "png", "jpg", "jpeg",
]

# Threshold para activar fallback OCR num PDF
OCR_TRIGGER_CHAR_THRESHOLD = 100

# Threshold de "texto sem sentido": ratio mínimo de chars alfa-numéricos
OCR_TRIGGER_GIBBERISH_RATIO = 0.30

# Tipo da callback de Vision LLM (injectada pelo app.py)
VisionCallback = Callable[[bytes], str]


# ─────────────────────────────────────────────────────────────────────────────
# DETECÇÃO DE TIPO E FALLBACK GIBBERISH
# ─────────────────────────────────────────────────────────────────────────────

def _detect_extension(file_name: str) -> str:
    """Devolve a extensão lowercase sem ponto, ou '' se sem extensão."""
    if not file_name or "." not in file_name:
        return ""
    return file_name.rsplit(".", 1)[-1].lower().strip()


def _looks_like_gibberish(text: str) -> bool:
    """
    Heurística para detectar texto sem sentido (típico de PDFs digitalizados
    com encoding partido onde o PyMuPDF devolve sequências de chars não-texto).

    Critério: ratio de chars alfa-numéricos < OCR_TRIGGER_GIBBERISH_RATIO.
    """
    if not text:
        return True
    sample = text[:5000]   # avalia apenas os primeiros 5k chars (rápido)
    if not sample.strip():
        return True
    alnum = sum(1 for c in sample if c.isalnum() or c.isspace())
    ratio = alnum / len(sample)
    return ratio < OCR_TRIGGER_GIBBERISH_RATIO


def _needs_ocr_fallback(text: str) -> Tuple[bool, str]:
    """
    Decide se deve ser activado o fallback OCR para um PDF.
    Devolve (precisa_ocr, motivo).
    """
    if not text or not text.strip():
        return True, "texto extraído vazio"
    if len(text.strip()) < OCR_TRIGGER_CHAR_THRESHOLD:
        return True, f"texto < {OCR_TRIGGER_CHAR_THRESHOLD} chars"
    if _looks_like_gibberish(text):
        return True, "texto parece sem sentido (ratio alfa-numérica baixa)"
    return False, ""


# ─────────────────────────────────────────────────────────────────────────────
# NORMALIZAÇÃO DE INPUT — aceita UploadedFile, BytesIO, bytes ou path
# ─────────────────────────────────────────────────────────────────────────────

def _read_file_bytes(file: Any) -> Tuple[bytes, str]:
    """
    Devolve (bytes, file_name) a partir de várias formas de input.

    Aceita:
      • Streamlit UploadedFile (tem .name e .getvalue())
      • bytes / bytearray (devolve sem nome — caller deve fornecer)
      • str ou Path (interpretado como ficheiro local)
      • io.BytesIO (precisa de .name atribuído manualmente)
    """
    # UploadedFile do Streamlit
    if hasattr(file, "getvalue") and hasattr(file, "name"):
        return file.getvalue(), file.name

    # Path / string
    if isinstance(file, str):
        from pathlib import Path
        path = Path(file)
        return path.read_bytes(), path.name

    # bytes puros
    if isinstance(file, (bytes, bytearray)):
        return bytes(file), getattr(file, "name", "unknown")

    # BytesIO
    if isinstance(file, io.BytesIO):
        data = file.getvalue()
        name = getattr(file, "name", "unknown")
        return data, name

    raise TypeError(
        f"Tipo de input não suportado para ingest_file: {type(file).__name__}. "
        "Esperado: UploadedFile, bytes, str (path) ou BytesIO."
    )


# ─────────────────────────────────────────────────────────────────────────────
# API PÚBLICA — ingest_file
# ─────────────────────────────────────────────────────────────────────────────

def ingest_file(
    file: Any,
    *,
    vision_callback: Optional[VisionCallback] = None,
) -> Dict[str, Any]:
    """
    Ingere um ficheiro de QUALQUER formato suportado e devolve um dict
    com texto consolidado, metadata e sub-textos por fonte.

    Args:
        file: UploadedFile (Streamlit), bytes, BytesIO, ou path string.
        vision_callback: opcional. Função (bytes) -> str que processa
                         imagens via Vision LLM. Se None, OCR cai para
                         Tesseract local ou skip com aviso.

    Returns:
        {
            "text":     str,                # texto consolidado pronto para LLM
            "metadata": {
                "file_name":          str,
                "file_type":          str,    # extensão lowercase
                "file_hash":          str,    # SHA-256
                "file_size_bytes":    int,
                "page_count":         int | None,
                "had_ocr_fallback":   bool,   # PDF caiu para OCR?
                "ingestion_warnings": list[str],
            },
            "sources": {
                "raw":    str,   # texto extraído digitalmente (PyMuPDF, etc.)
                "ocr":    str,   # texto via OCR (PDF scanned ou imagens)
                "tables": str,   # tabelas de xlsx/docx
            },
        }

    Nunca levanta exception por formato não suportado ou parser falhado —
    devolve sempre o dict com warnings em metadata.ingestion_warnings.
    """
    try:
        file_bytes, file_name = _read_file_bytes(file)
    except TypeError as e:
        return _empty_result(file_name="unknown", file_type="",
                             warnings=[f"input inválido: {e}"])

    file_type = _detect_extension(file_name)
    file_hash = compute_file_hash(file_bytes)
    file_size = len(file_bytes)

    metadata: Dict[str, Any] = {
        "file_name": file_name,
        "file_type": file_type,
        "file_hash": file_hash,
        "file_size_bytes": file_size,
        "page_count": None,
        "had_ocr_fallback": False,
        "ingestion_warnings": [],
    }
    sources: Dict[str, str] = {"raw": "", "ocr": "", "tables": ""}

    if file_type not in SUPPORTED_FORMATS:
        metadata["ingestion_warnings"].append(
            f"Formato '{file_type}' não suportado. "
            f"Suportados: {', '.join(SUPPORTED_FORMATS)}."
        )
        return {"text": "", "metadata": metadata, "sources": sources}

    # ─── Roteamento por formato ──────────────────────────────────────────
    try:
        if file_type == "pdf":
            sources, metadata = _parse_pdf(
                file_bytes, sources, metadata, vision_callback=vision_callback,
            )

        elif file_type == "docx":
            sources, metadata = _parse_docx(file_bytes, sources, metadata)

        elif file_type == "pptx":
            sources, metadata = _parse_pptx(file_bytes, sources, metadata)

        elif file_type == "xlsx":
            sources, metadata = _parse_xlsx(file_bytes, sources, metadata)

        elif file_type in ("txt", "md"):
            sources, metadata = _parse_plain_text(file_bytes, sources, metadata)

        elif file_type == "html":
            sources, metadata = _parse_html(file_bytes, sources, metadata)

        elif file_type in ("png", "jpg", "jpeg"):
            sources, metadata = _parse_image(
                file_bytes, sources, metadata, vision_callback=vision_callback,
            )

    except Exception as e:
        metadata["ingestion_warnings"].append(
            f"Erro inesperado ao processar '{file_name}': {e}"
        )

    # Texto consolidado: concatena fontes não-vazias por ordem natural
    consolidated_parts = []
    if sources["raw"].strip():
        consolidated_parts.append(sources["raw"].strip())
    if sources["ocr"].strip() and sources["ocr"].strip() != sources["raw"].strip():
        consolidated_parts.append("[--- OCR ---]\n" + sources["ocr"].strip())
    if sources["tables"].strip():
        consolidated_parts.append("[--- TABLES ---]\n" + sources["tables"].strip())

    text = "\n\n".join(consolidated_parts)

    return {"text": text, "metadata": metadata, "sources": sources}


# ─────────────────────────────────────────────────────────────────────────────
# Parser: PDF (com regra crítica de fallback)
# ─────────────────────────────────────────────────────────────────────────────

def _parse_pdf(
    file_bytes: bytes,
    sources: Dict[str, str],
    metadata: Dict[str, Any],
    *,
    vision_callback: Optional[VisionCallback],
) -> Tuple[Dict[str, str], Dict[str, Any]]:
    """
    PDF: tenta PyMuPDF primeiro. Se texto < 100 chars OU detecção de
    gibberish, accionar fallback OCR (Vision LLM ou Tesseract).
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        metadata["ingestion_warnings"].append(
            "PyMuPDF não está instalado. `pip install PyMuPDF`."
        )
        return sources, metadata

    raw_text = ""
    n_pages = 0
    try:
        with fitz.open(stream=file_bytes, filetype="pdf") as doc:
            n_pages = len(doc)
            page_texts = [page.get_text() for page in doc]
            raw_text = "\n\n".join(page_texts).strip()
    except Exception as e:
        metadata["ingestion_warnings"].append(f"Falha PyMuPDF: {e}")
        raw_text = ""

    metadata["page_count"] = n_pages
    sources["raw"] = raw_text

    # ─── Regra crítica de fallback (spec §5) ───────────────────────────────
    needs_ocr, reason = _needs_ocr_fallback(raw_text)
    if not needs_ocr:
        return sources, metadata

    metadata["ingestion_warnings"].append(
        f"PDF accionou fallback OCR ({reason})."
    )

    ocr_text = _run_pdf_ocr(file_bytes, vision_callback=vision_callback,
                            warnings=metadata["ingestion_warnings"])
    sources["ocr"] = ocr_text
    metadata["had_ocr_fallback"] = True
    return sources, metadata


def _run_pdf_ocr(
    file_bytes: bytes,
    *,
    vision_callback: Optional[VisionCallback],
    warnings: List[str],
) -> str:
    """
    Cadeia de fallback OCR:
      1. Vision LLM via vision_callback (preferido se passado)
      2. Tesseract + pdf2image (se ambos instalados localmente)
      3. Devolve "" + warning
    """
    # ── Caminho 1: Vision LLM (qualidade superior) ─────────────────────────
    if vision_callback is not None:
        try:
            return _ocr_via_vision_llm(file_bytes, vision_callback, warnings)
        except Exception as e:
            warnings.append(f"Vision LLM falhou ({e}); tentando Tesseract local…")

    # ── Caminho 2: Tesseract + pdf2image ──────────────────────────────────
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
    except ImportError:
        warnings.append(
            "OCR local indisponível: instala `pdf2image` + `pytesseract` "
            "(e tem `poppler` + `tesseract` no sistema) — OU passa um "
            "vision_callback (Vision LLM)."
        )
        return ""

    try:
        images = convert_from_bytes(file_bytes, dpi=200)
        pages_text = []
        for i, img in enumerate(images, start=1):
            try:
                pages_text.append(pytesseract.image_to_string(img))
            except Exception as e:
                warnings.append(f"Tesseract falhou na página {i}: {e}")
        return "\n\n".join(t for t in pages_text if t.strip())
    except Exception as e:
        warnings.append(f"OCR via Tesseract falhou: {e}")
        return ""


def _ocr_via_vision_llm(
    file_bytes: bytes,
    vision_callback: VisionCallback,
    warnings: List[str],
) -> str:
    """
    Converte cada página do PDF em imagem PNG e invoca o vision_callback
    para extrair texto. Concatena os resultados.
    """
    try:
        from pdf2image import convert_from_bytes
    except ImportError:
        # Fallback: tenta extrair imagens via PyMuPDF directamente
        return _ocr_via_vision_pymupdf(file_bytes, vision_callback, warnings)

    pages_text = []
    try:
        images = convert_from_bytes(file_bytes, dpi=200)
        for i, pil_img in enumerate(images, start=1):
            buf = io.BytesIO()
            pil_img.save(buf, format="PNG")
            try:
                page_text = vision_callback(buf.getvalue())
                if page_text:
                    pages_text.append(page_text)
            except Exception as e:
                warnings.append(f"Vision LLM falhou na página {i}: {e}")
    except Exception as e:
        warnings.append(f"Falha a converter PDF→imagens: {e}")
    return "\n\n".join(pages_text)


def _ocr_via_vision_pymupdf(
    file_bytes: bytes,
    vision_callback: VisionCallback,
    warnings: List[str],
) -> str:
    """
    Fallback se `pdf2image` não estiver instalado: usa PyMuPDF para
    renderizar páginas como pixmap e converter para PNG.
    """
    try:
        import fitz
    except ImportError:
        warnings.append("PyMuPDF não disponível para fallback de Vision.")
        return ""

    pages_text = []
    try:
        with fitz.open(stream=file_bytes, filetype="pdf") as doc:
            for i, page in enumerate(doc, start=1):
                try:
                    pix = page.get_pixmap(dpi=200)
                    png_bytes = pix.tobytes("png")
                    page_text = vision_callback(png_bytes)
                    if page_text:
                        pages_text.append(page_text)
                except Exception as e:
                    warnings.append(f"Vision LLM falhou na página {i}: {e}")
    except Exception as e:
        warnings.append(f"PyMuPDF→imagens falhou: {e}")
    return "\n\n".join(pages_text)


# ─────────────────────────────────────────────────────────────────────────────
# Parser: DOCX
# ─────────────────────────────────────────────────────────────────────────────

def _parse_docx(
    file_bytes: bytes,
    sources: Dict[str, str],
    metadata: Dict[str, Any],
) -> Tuple[Dict[str, str], Dict[str, Any]]:
    try:
        from docx import Document
    except ImportError:
        metadata["ingestion_warnings"].append(
            "python-docx não instalado. `pip install python-docx`."
        )
        return sources, metadata

    try:
        doc = Document(io.BytesIO(file_bytes))
    except Exception as e:
        metadata["ingestion_warnings"].append(f"Falha a abrir DOCX: {e}")
        return sources, metadata

    # Parágrafos
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    # Tabelas
    table_chunks = []
    for ti, table in enumerate(doc.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            table_chunks.append(f"### Tabela {ti}\n" + "\n".join(rows))

    sources["raw"] = "\n\n".join(paragraphs)
    sources["tables"] = "\n\n".join(table_chunks)
    metadata["page_count"] = None   # DOCX não tem páginas estáveis
    return sources, metadata


# ─────────────────────────────────────────────────────────────────────────────
# Parser: PPTX
# ─────────────────────────────────────────────────────────────────────────────

def _parse_pptx(
    file_bytes: bytes,
    sources: Dict[str, str],
    metadata: Dict[str, Any],
) -> Tuple[Dict[str, str], Dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError:
        metadata["ingestion_warnings"].append(
            "python-pptx não instalado. `pip install python-pptx`."
        )
        return sources, metadata

    try:
        pres = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        metadata["ingestion_warnings"].append(f"Falha a abrir PPTX: {e}")
        return sources, metadata

    slides_text = []
    n_slides = 0
    for i, slide in enumerate(pres.slides, start=1):
        n_slides = i
        parts = [f"## Slide {i}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            # Notas do orador
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"_[Notas do orador]_: {notes}")
        slides_text.append("\n\n".join(parts))

    sources["raw"] = "\n\n---\n\n".join(slides_text)
    metadata["page_count"] = n_slides
    return sources, metadata


# ─────────────────────────────────────────────────────────────────────────────
# Parser: XLSX
# ─────────────────────────────────────────────────────────────────────────────

def _parse_xlsx(
    file_bytes: bytes,
    sources: Dict[str, str],
    metadata: Dict[str, Any],
) -> Tuple[Dict[str, str], Dict[str, Any]]:
    """
    Lê XLSX em texto Markdown-friendly. Tenta pandas (mais limpo); se
    pandas não estiver instalado, cai para openpyxl directo.
    """
    # ── Tentativa 1: pandas ──
    try:
        import pandas as pd
        try:
            sheets = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None, header=None)
            chunks = []
            for sheet_name, df in sheets.items():
                # Cap em 200 linhas para evitar prompts gigantes
                df_sample = df.head(200)
                chunks.append(
                    f"## Folha: {sheet_name}\n\n```\n{df_sample.to_string(index=False, header=False)}\n```"
                )
            sources["tables"] = "\n\n".join(chunks)
            metadata["page_count"] = len(sheets)
            return sources, metadata
        except Exception as e:
            metadata["ingestion_warnings"].append(f"pandas falhou XLSX: {e}; a tentar openpyxl direto…")
    except ImportError:
        pass   # pandas não instalado; cai para openpyxl

    # ── Tentativa 2: openpyxl direct ──
    try:
        from openpyxl import load_workbook
    except ImportError:
        metadata["ingestion_warnings"].append(
            "Nem pandas nem openpyxl instalados. `pip install openpyxl`."
        )
        return sources, metadata

    try:
        wb = load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
        chunks = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 200:
                    rows.append("[... truncado a 200 linhas ...]")
                    break
                cells = ["" if c is None else str(c) for c in row]
                rows.append(" | ".join(cells))
            chunks.append(f"## Folha: {sheet_name}\n\n```\n" + "\n".join(rows) + "\n```")
        wb.close()
        sources["tables"] = "\n\n".join(chunks)
        metadata["page_count"] = len(wb.sheetnames)
    except Exception as e:
        metadata["ingestion_warnings"].append(f"openpyxl falhou XLSX: {e}")
    return sources, metadata


# ─────────────────────────────────────────────────────────────────────────────
# Parser: TXT / MD (leitura directa)
# ─────────────────────────────────────────────────────────────────────────────

def _parse_plain_text(
    file_bytes: bytes,
    sources: Dict[str, str],
    metadata: Dict[str, Any],
) -> Tuple[Dict[str, str], Dict[str, Any]]:
    text = _decode_text_bytes(file_bytes, metadata)
    sources["raw"] = text
    return sources, metadata


# ─────────────────────────────────────────────────────────────────────────────
# Parser: HTML (strip tags básico, sem dependências)
# ─────────────────────────────────────────────────────────────────────────────

_HTML_TAG_RE = re.compile(r"<[^>]+>")
_HTML_SCRIPT_STYLE_RE = re.compile(
    r"<(script|style)[^>]*>.*?</\1>", re.DOTALL | re.IGNORECASE,
)
_HTML_ENTITIES = {"&nbsp;": " ", "&amp;": "&", "&lt;": "<", "&gt;": ">",
                  "&quot;": '"', "&apos;": "'", "&#39;": "'", "&hellip;": "…"}


def _parse_html(
    file_bytes: bytes,
    sources: Dict[str, str],
    metadata: Dict[str, Any],
) -> Tuple[Dict[str, str], Dict[str, Any]]:
    """
    Strip de tags HTML básico. Para HTMLs complexos, recomendar
    BeautifulSoup4 (mas não exigir).
    """
    raw_html = _decode_text_bytes(file_bytes, metadata)

    # Tenta BeautifulSoup se disponível (melhor qualidade)
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw_html, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        # Colapsa linhas vazias múltiplas
        text = re.sub(r"\n{3,}", "\n\n", text)
        sources["raw"] = text
        return sources, metadata
    except ImportError:
        pass   # cai para regex

    # Fallback regex (sem dependências externas)
    text = _HTML_SCRIPT_STYLE_RE.sub("", raw_html)
    text = _HTML_TAG_RE.sub("\n", text)
    for entity, replacement in _HTML_ENTITIES.items():
        text = text.replace(entity, replacement)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    sources["raw"] = text.strip()
    return sources, metadata


# ─────────────────────────────────────────────────────────────────────────────
# Parser: Imagem (PNG/JPG/JPEG) — só funciona com vision_callback
# ─────────────────────────────────────────────────────────────────────────────

def _parse_image(
    file_bytes: bytes,
    sources: Dict[str, str],
    metadata: Dict[str, Any],
    *,
    vision_callback: Optional[VisionCallback],
) -> Tuple[Dict[str, str], Dict[str, Any]]:
    """
    PNG/JPG: tem que ir pelo pipeline de Vision (LLM multimodal).
    Se vision_callback não foi passada, tenta Tesseract local; senão skip.
    """
    if vision_callback is not None:
        try:
            text = vision_callback(file_bytes)
            sources["ocr"] = text or ""
            metadata["had_ocr_fallback"] = True
            return sources, metadata
        except Exception as e:
            metadata["ingestion_warnings"].append(
                f"Vision LLM falhou na imagem: {e}; a tentar Tesseract…"
            )

    # Fallback Tesseract para imagem solta
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(img)
        sources["ocr"] = text or ""
        metadata["had_ocr_fallback"] = True
    except ImportError:
        metadata["ingestion_warnings"].append(
            "Imagem requer vision_callback (Vision LLM) OU `pytesseract + Pillow` "
            "(com `tesseract` instalado no sistema)."
        )
    except Exception as e:
        metadata["ingestion_warnings"].append(f"Tesseract falhou na imagem: {e}")
    return sources, metadata


# ─────────────────────────────────────────────────────────────────────────────
# Utilidades de decoding
# ─────────────────────────────────────────────────────────────────────────────

def _decode_text_bytes(file_bytes: bytes, metadata: Dict[str, Any]) -> str:
    """
    Decode robusto: tenta utf-8, depois utf-8-sig, latin-1 como último recurso.
    Aviso registado em metadata se forem necessários encodings de fallback.
    """
    for encoding in ("utf-8", "utf-8-sig"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    metadata["ingestion_warnings"].append(
        "Encoding não é UTF-8; a usar latin-1 com possíveis substituições."
    )
    return file_bytes.decode("latin-1", errors="replace")


# ─────────────────────────────────────────────────────────────────────────────
# Resultado vazio (helper para erro inicial)
# ─────────────────────────────────────────────────────────────────────────────

def _empty_result(
    file_name: str,
    file_type: str,
    warnings: List[str],
) -> Dict[str, Any]:
    return {
        "text": "",
        "metadata": {
            "file_name": file_name,
            "file_type": file_type,
            "file_hash": "",
            "file_size_bytes": 0,
            "page_count": None,
            "had_ocr_fallback": False,
            "ingestion_warnings": warnings,
        },
        "sources": {"raw": "", "ocr": "", "tables": ""},
    }
